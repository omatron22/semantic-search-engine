import sys
from pptx import Presentation


def extract_text_from_pptx(pptx_path):
    """Extract text from PowerPoint file by reading all slides."""
    try:
        prs = Presentation(pptx_path)
        lines = []

        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"Slide {i}:")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            lines.append(text)

        return "\n".join(lines).strip()
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python parse_pptx.py <path_to_pptx>")
        sys.exit(1)

    text = extract_text_from_pptx(sys.argv[1])
    print(text)
